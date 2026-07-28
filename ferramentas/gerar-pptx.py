#!/usr/bin/env python3
"""
Gera o .pptx de entrega a partir de um slides.html.

    python3 ferramentas/gerar-pptx.py palestras/<nome>/slides.html

Renderiza cada slide no MODO LIMPO (sem a coluna do formador), em 1920x1080, e
monta um PowerPoint 16:9 com cada slide como imagem de página inteira. As notas
do formador vão para o campo de anotações do PowerPoint — então o arquivo
entregue continua sendo o "PPT do palestrante" quando aberto no modo
apresentador, e o "PPT final" quando projetado.

Por que imagem e não formas nativas: o layout usa container queries, cqw e
fontes próprias. Reconstruir isso em formas do PowerPoint perderia exatamente o
acabamento que a identidade exige. A imagem preserva o slide como ele é.

Dependências (instalar uma vez):
    pip install playwright python-pptx
    playwright install chromium      # pule se PLAYWRIGHT_BROWSERS_PATH já aponta para um chromium
"""

import asyncio
import os
import pathlib
import sys

LARGURA, ALTURA = 1920, 1080


def erro(msg: str) -> None:
    print(f"erro: {msg}", file=sys.stderr)
    raise SystemExit(1)


def chromium_alternativo():
    """Procura um Chromium já instalado, para quando o build que o Playwright
    espera não bater com o que existe na máquina (acontece em contêiner com
    PLAYWRIGHT_BROWSERS_PATH pré-configurado)."""
    base = os.environ.get("PLAYWRIGHT_BROWSERS_PATH")
    if not base or not pathlib.Path(base).is_dir():
        return None
    candidatos = sorted(pathlib.Path(base).glob("chromium-*/chrome-linux/chrome"))
    return str(candidatos[-1]) if candidatos else None


async def abrir_navegador(p):
    try:
        return await p.chromium.launch()
    except Exception:
        alternativo = chromium_alternativo()
        if not alternativo:
            erro("Chromium não encontrado. Rode: playwright install chromium")
        print(f"usando chromium de {alternativo}")
        return await p.chromium.launch(executable_path=alternativo)


async def capturar(origem: pathlib.Path, destino: pathlib.Path):
    try:
        from playwright.async_api import async_playwright
    except ImportError:
        erro("playwright não instalado. Rode: pip install playwright python-pptx")

    destino.mkdir(parents=True, exist_ok=True)
    imagens, notas = [], []

    async with async_playwright() as p:
        navegador = await abrir_navegador(p)
        pagina = await navegador.new_page(
            viewport={"width": LARGURA, "height": ALTURA}, device_scale_factor=2
        )
        await pagina.goto(origem.resolve().as_uri() + "?modo=limpo")
        await pagina.wait_for_function("window.__palestra !== undefined")
        await pagina.wait_for_timeout(900)  # fontes

        total = await pagina.evaluate("window.__palestra.total")
        notas = await pagina.evaluate("window.__palestra.notas()")
        print(f"{total} slides encontrados")

        for i in range(total):
            await pagina.evaluate(f"window.__palestra.irPara({i})")
            await pagina.wait_for_timeout(220)
            caminho = destino / f"slide-{i + 1:02d}.png"
            await pagina.locator(".slide.ativo .tela").screenshot(path=str(caminho))
            imagens.append(caminho)
            print(f"  slide {i + 1}/{total}")

        await navegador.close()

    return imagens, notas


def montar_pptx(imagens, notas, saida: pathlib.Path):
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        erro("python-pptx não instalado. Rode: pip install python-pptx")

    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 16:9
    prs.slide_height = Inches(7.5)
    em_branco = prs.slide_layouts[6]

    for i, img in enumerate(imagens):
        slide = prs.slides.add_slide(em_branco)
        slide.shapes.add_picture(
            str(img), 0, 0, width=prs.slide_width, height=prs.slide_height
        )
        if i < len(notas) and notas[i]:
            slide.notes_slide.notes_text_frame.text = notas[i]

    prs.save(str(saida))
    print(f"\n✓ {saida}  ({len(imagens)} slides, notas do formador incluídas)")


def main():
    if len(sys.argv) < 2:
        erro("uso: python3 ferramentas/gerar-pptx.py palestras/<nome>/slides.html")

    origem = pathlib.Path(sys.argv[1])
    if not origem.exists():
        erro(f"não encontrei {origem}")

    pasta = origem.parent
    temp = pasta / ".render"
    saida = pathlib.Path(sys.argv[2]) if len(sys.argv) > 2 else pasta / f"{pasta.name}.pptx"

    imagens, notas = asyncio.run(capturar(origem, temp))
    if not imagens:
        erro("nenhum slide capturado — o slides.html carrega o slides.js?")
    montar_pptx(imagens, notas, saida)
    print(f"  PNGs em {temp}/ (pode apagar, ou usar para divulgação)")


if __name__ == "__main__":
    main()
